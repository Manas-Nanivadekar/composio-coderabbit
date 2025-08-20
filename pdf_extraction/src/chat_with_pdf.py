from fastapi import FastAPI, UploadFile, File, HTTPException, Form
from fastapi.middleware.cors import CORSMiddleware
from typing import List, Dict, Optional
from PyPDF2 import PdfReader
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
import asyncio
import tempfile
import json
import uuid
from datetime import datetime
from pathlib import Path

from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai

from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv

# Load environment variables and configure Google API
load_dotenv()
if not os.getenv("GOOGLE_API_KEY"):
    raise RuntimeError("GOOGLE_API_KEY is not set in environment")
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

app = FastAPI(title="PDF Chat API", version="1.0.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Global storage for sessions and documents
sessions: Dict[str, Dict] = {}
SESSION_DIR = Path("sessions")
SESSION_DIR.mkdir(exist_ok=True)

# Function to extract text from PDF files
def get_pdf_text(pdf_docs: List[UploadFile]) -> str:
    text = ""
    for pdf in pdf_docs:
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                data = pdf.file.read()
                tmp.write(data)
                tmp_path = tmp.name
            pdf_reader = PdfReader(tmp_path)
            for page in pdf_reader.pages:
                text += (page.extract_text() or "")
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Failed to read PDF '{pdf.filename}': {e}")
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.remove(tmp_path)
                except Exception:
                    pass
    return text

# Function to extract text from PPTX files (kept for compatibility)
def get_pptx_text(pptx_docs: List[UploadFile]) -> str:
    text = ""
    for pptx in pptx_docs:
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                data = pptx.file.read()
                tmp.write(data)
                tmp_path = tmp.name
            presentation = Presentation(tmp_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Failed to read PPTX '{pptx.filename}': {e}")
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.remove(tmp_path)
                except Exception:
                    pass
    return text

# Function to split text into chunks
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

# Function to create and save vector store with session support
def get_vector_store(text_chunks, session_id: str):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    
    # Create session-specific directory
    session_path = SESSION_DIR / session_id
    session_path.mkdir(exist_ok=True)
    
    # Save vector store with session ID
    vector_store.save_local(str(session_path / "faiss_index"))
    return vector_store

# Function to create conversational chain with improved prompt
def get_conversational_chain():
    prompt_template = (
        "You are a helpful assistant that answers questions based on the provided document context. "
        "Use the context below to answer the question as accurately and comprehensively as possible. "
        "If the exact answer is not in the context, provide the most relevant information available "
        "and indicate what information might be missing.\n\n"
        "Context:\n{context}\n\n"
        "Question: {question}\n\n"
        "Answer (be specific and helpful):"
    )
    model = ChatGoogleGenerativeAI(model="gemini-2.0-flash", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

# Function to handle user input with session support and improved retrieval
def answer_question(user_question: str, session_id: str) -> str:
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    session_path = SESSION_DIR / session_id / "faiss_index"
    
    try:
        new_db = FAISS.load_local(str(session_path), embeddings, allow_dangerous_deserialization=True)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Vector index not found for session. Please upload documents first. Error: {e}")
    
    # Improved similarity search with more results and score threshold
    docs = new_db.similarity_search(user_question, k=5)
    
    # Debug: Print retrieved documents (remove in production)
    print(f"Found {len(docs)} documents for question: {user_question}")
    for i, doc in enumerate(docs):
        print(f"Doc {i+1} preview: {doc.page_content[:200]}...")
    
    if not docs:
        return "I couldn't find any relevant information in the uploaded documents. Please try rephrasing your question or upload relevant documents."
    
    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    return response["output_text"]

@app.post("/extract_pdf")
async def extract_pdf(files: List[UploadFile] = File(...), session_id: str = Form(None)):
    # Generate session ID if not provided
    if not session_id:
        session_id = str(uuid.uuid4())
    
    pdf_files = [f for f in files if f.filename.lower().endswith(".pdf")]
    if not pdf_files:
        raise HTTPException(status_code=400, detail="No PDF files uploaded")

    # Extract text only from PDFs
    raw_text = get_pdf_text(pdf_files)
    if not raw_text.strip():
        raise HTTPException(status_code=400, detail="No extractable text found in uploaded PDFs")

    text_chunks = get_text_chunks(raw_text)
    vector_store = get_vector_store(text_chunks, session_id)
    
    # Store session information with simplified file info
    file_info = []
    for f in pdf_files:
        file_info.append({
            "name": f.filename,
            "size": "Unknown",  # Simplified - size isn't critical for functionality
            "upload_date": datetime.now().isoformat()
        })
    
    sessions[session_id] = {
        "created_at": datetime.now().isoformat(),
        "documents": file_info,
        "chunk_count": len(text_chunks),
        "text_length": len(raw_text)
    }
    
    # Save session to disk for persistence
    session_path = SESSION_DIR / session_id
    session_path.mkdir(exist_ok=True)
    with open(session_path / "session.json", "w") as f:
        json.dump(sessions[session_id], f, indent=2)
    
    return {
        "message": "Index built successfully", 
        "chunk_count": len(text_chunks), 
        "session_id": session_id,
        "documents": sessions[session_id]["documents"]
    }


@app.post("/chat")
async def chat(question: str = Form(...), session_id: str = Form(...)):
    if not question or not question.strip():
        raise HTTPException(status_code=400, detail="Question must not be empty")
    if not session_id:
        raise HTTPException(status_code=400, detail="Session ID is required")
    
    answer = answer_question(question, session_id)
    return {"answer": answer}

@app.get("/sessions/{session_id}")
async def get_session(session_id: str):
    """Get session information and uploaded documents"""
    session_path = SESSION_DIR / session_id
    session_file = session_path / "session.json"
    
    if not session_file.exists():
        raise HTTPException(status_code=404, detail="Session not found")
    
    # Load session from disk if not in memory
    if session_id not in sessions:
        with open(session_file, "r") as f:
            sessions[session_id] = json.load(f)
    
    return sessions[session_id]

@app.get("/sessions")
async def list_sessions():
    """List all available sessions"""
    session_dirs = [d for d in SESSION_DIR.iterdir() if d.is_dir()]
    session_list = []
    
    for session_dir in session_dirs:
        session_file = session_dir / "session.json"
        if session_file.exists():
            with open(session_file, "r") as f:
                session_data = json.load(f)
                session_list.append({
                    "session_id": session_dir.name,
                    "created_at": session_data.get("created_at"),
                    "document_count": len(session_data.get("documents", [])),
                    "chunk_count": session_data.get("chunk_count", 0)
                })
    
    return {"sessions": session_list}

@app.delete("/sessions/{session_id}")
async def delete_session(session_id: str):
    """Delete a session and its associated data"""
    session_path = SESSION_DIR / session_id
    
    if not session_path.exists():
        raise HTTPException(status_code=404, detail="Session not found")
    
    # Remove from memory
    if session_id in sessions:
        del sessions[session_id]
    
    # Remove files
    import shutil
    shutil.rmtree(session_path)
    
    return {"message": "Session deleted successfully"}
